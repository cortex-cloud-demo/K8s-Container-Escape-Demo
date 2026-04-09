#!/usr/bin/env python3
"""Generate a PowerPoint presentation for K8s Container Escape Demo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x81, 0x8C, 0xF8)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)
DARK_BG2 = RGBColor(0x1A, 0x23, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=DARK_CARD, border_color=None, border_width=Pt(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=WHITE, bullet_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = bullet_color or color
        p.space_before = Pt(6)
        p.level = 0
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Kubernetes Container Escape Demo", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "From Spring4Shell RCE to Cluster Takeover", font_size=24, color=BLUE, bold=False, alignment=PP_ALIGN.CENTER)

# Description
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(1.2),
             "Full attack chain on AWS EKS with automated detection by Cortex XDR\nand incident response via Cortex Playbooks + AWS Lambda containment",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(3), Inches(5.8), Inches(7), Inches(0.6), fill_color=DARK_CARD, border_color=BLUE)
add_text_box(slide, Inches(3), Inches(5.85), Inches(7), Inches(0.5),
             "Palo Alto Networks  |  Cortex Cloud Demo", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Architecture Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Architecture Overview", font_size=32, color=WHITE, bold=True)

# Add architecture image
arch_img = os.path.join(os.path.dirname(__file__), "architecture.png")
if os.path.exists(arch_img):
    slide.shapes.add_picture(arch_img, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.8))
else:
    # Fallback: draw boxes
    # Operator box
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(3.5), Inches(2.5), border_color=BLUE)
    add_text_box(slide, Inches(0.7), Inches(1.4), Inches(3), Inches(0.4), "OPERATOR", font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(0.7), Inches(1.9), Inches(3), Inches(0.4), "Web Dashboard", font_size=13, color=WHITE)
    add_text_box(slide, Inches(0.7), Inches(2.3), Inches(3), Inches(0.4), "Cortex XSIAM", font_size=13, color=PURPLE)

    # AWS box
    add_shape(slide, Inches(5), Inches(1.3), Inches(7.5), Inches(5.5), border_color=ORANGE)
    add_text_box(slide, Inches(5.2), Inches(1.4), Inches(3), Inches(0.4), "AWS CLOUD", font_size=14, color=ORANGE, bold=True)

    # EKS
    add_shape(slide, Inches(5.5), Inches(2.2), Inches(4.5), Inches(3.5), border_color=GREEN)
    add_text_box(slide, Inches(5.7), Inches(2.3), Inches(4), Inches(0.4), "Amazon EKS", font_size=13, color=GREEN, bold=True)
    add_text_box(slide, Inches(5.7), Inches(2.8), Inches(4), Inches(2.5),
                 "Namespace: vuln-app\n"
                 "  - Privileged Pod (Spring4Shell)\n"
                 "  - SA cluster-admin\n"
                 "  - LoadBalancer\n\n"
                 "Node Group: 2x t3.medium (AL2023)",
                 font_size=11, color=GRAY)

    # Lambda
    add_shape(slide, Inches(10.5), Inches(2.2), Inches(2), Inches(1.5), border_color=ORANGE)
    add_text_box(slide, Inches(10.6), Inches(2.3), Inches(1.8), Inches(1.3),
                 "Lambda\nContainment\n7 actions", font_size=11, color=ORANGE)


# ============================================================
# SLIDE 3: Attack Chain
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Attack Chain — 3 Steps", font_size=32, color=RED, bold=True)

# Step 1
add_shape(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(5.5), border_color=RED)
add_text_box(slide, Inches(0.7), Inches(1.4), Inches(3.4), Inches(0.5),
             "STEP 1: Spring4Shell RCE", font_size=18, color=RED, bold=True)
txBox = add_text_box(slide, Inches(0.7), Inches(2.0), Inches(3.4), Inches(4.5), "", font_size=13, color=WHITE)
tf = txBox.text_frame
tf.paragraphs[0].text = "CVE-2022-22965 (CVSS 9.8)"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = ORANGE
tf.paragraphs[0].font.bold = True
add_paragraph(tf, "")
add_paragraph(tf, "1. Exploit ClassLoader via Spring", font_size=12, color=GRAY)
add_paragraph(tf, "   parameter binding", font_size=12, color=GRAY)
add_paragraph(tf, "2. Manipulate Tomcat", font_size=12, color=GRAY)
add_paragraph(tf, "   AccessLogValve", font_size=12, color=GRAY)
add_paragraph(tf, "3. Write JSP webshell to", font_size=12, color=GRAY)
add_paragraph(tf, "   webapps/app/", font_size=12, color=GRAY)
add_paragraph(tf, "4. Verify RCE: id command", font_size=12, color=GRAY)
add_paragraph(tf, "")
add_paragraph(tf, "Result:", font_size=13, color=WHITE, bold=True)
add_paragraph(tf, "Remote Code Execution", font_size=13, color=RED)
add_paragraph(tf, "on the container", font_size=13, color=RED)

# Step 2
add_shape(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(5.5), border_color=RED)
add_text_box(slide, Inches(4.9), Inches(1.4), Inches(3.4), Inches(0.5),
             "STEP 2: Container Escape", font_size=18, color=RED, bold=True)
txBox = add_text_box(slide, Inches(4.9), Inches(2.0), Inches(3.4), Inches(4.5), "", font_size=13, color=WHITE)
tf = txBox.text_frame
tf.paragraphs[0].text = "Privileged + hostPID + hostPath"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = ORANGE
tf.paragraphs[0].font.bold = True
add_paragraph(tf, "")
add_paragraph(tf, "1. nsenter -t 1 -m -u -i -n -p", font_size=12, color=GRAY)
add_paragraph(tf, "   → escape to host namespace", font_size=12, color=GRAY)
add_paragraph(tf, "2. Read /etc/shadow, /etc/passwd", font_size=12, color=GRAY)
add_paragraph(tf, "3. Access host filesystem via", font_size=12, color=GRAY)
add_paragraph(tf, "   /proc/1/root", font_size=12, color=GRAY)
add_paragraph(tf, "4. IMDS 169.254.169.254", font_size=12, color=GRAY)
add_paragraph(tf, "   → steal AWS credentials", font_size=12, color=GRAY)
add_paragraph(tf, "")
add_paragraph(tf, "Result:", font_size=13, color=WHITE, bold=True)
add_paragraph(tf, "Full node access +", font_size=13, color=RED)
add_paragraph(tf, "AWS IAM credentials", font_size=13, color=RED)

# Step 3
add_shape(slide, Inches(8.9), Inches(1.3), Inches(3.8), Inches(5.5), border_color=RED)
add_text_box(slide, Inches(9.1), Inches(1.4), Inches(3.4), Inches(0.5),
             "STEP 3: Cluster Takeover", font_size=18, color=RED, bold=True)
txBox = add_text_box(slide, Inches(9.1), Inches(2.0), Inches(3.4), Inches(4.5), "", font_size=13, color=WHITE)
tf = txBox.text_frame
tf.paragraphs[0].text = "SA cluster-admin token"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = ORANGE
tf.paragraphs[0].font.bold = True
add_paragraph(tf, "")
add_paragraph(tf, "1. Steal ServiceAccount token", font_size=12, color=GRAY)
add_paragraph(tf, "   from /var/run/secrets/", font_size=12, color=GRAY)
add_paragraph(tf, "2. List all K8s secrets", font_size=12, color=GRAY)
add_paragraph(tf, "   (all namespaces)", font_size=12, color=GRAY)
add_paragraph(tf, "3. Extract AWS credentials", font_size=12, color=GRAY)
add_paragraph(tf, "   from secrets", font_size=12, color=GRAY)
add_paragraph(tf, "4. Full cluster control", font_size=12, color=GRAY)
add_paragraph(tf, "")
add_paragraph(tf, "Result:", font_size=13, color=WHITE, bold=True)
add_paragraph(tf, "Complete cluster +", font_size=13, color=RED)
add_paragraph(tf, "AWS compromise", font_size=13, color=RED)

# Arrows between steps
for x in [Inches(4.35), Inches(8.55)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(3.8), Inches(0.35), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RED
    arrow.line.fill.background()


# ============================================================
# SLIDE 4: Misconfigurations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Pod Misconfigurations Exploited", font_size=32, color=ORANGE, bold=True)

misconfigs = [
    ("privileged: true", "Full host kernel access", "allowPrivilegeEscalation: false", RED),
    ("hostPID: true", "Host process visibility, nsenter escape", "Disable hostPID", RED),
    ("hostNetwork: true", "Node network access, IMDS", "Disable, IMDSv2 hop limit=1", RED),
    ("hostPath: /", "Read/write entire host filesystem", "PVCs, restrict via PSA", RED),
    ("SA cluster-admin", "Full K8s API control", "Least privilege RBAC", RED),
    ("No Pod Security Standards", "All misconfigs allowed", "Enforce restricted PSA", ORANGE),
    ("No Network Policies", "Unrestricted pod communication", "Implement NetworkPolicies", ORANGE),
]

headers = ["Misconfiguration", "Impact", "Remediation"]
col_x = [Inches(0.5), Inches(4.0), Inches(8.5)]
col_w = [Inches(3.3), Inches(4.3), Inches(4.0)]

# Header row
add_shape(slide, Inches(0.4), Inches(1.2), Inches(12.4), Inches(0.5), fill_color=RGBColor(0x1E, 0x3A, 0x5F), border_color=BLUE)
for i, h in enumerate(headers):
    add_text_box(slide, col_x[i], Inches(1.25), col_w[i], Inches(0.4), h, font_size=14, color=BLUE, bold=True)

# Data rows
for row_i, (misconf, impact, fix, color) in enumerate(misconfigs):
    y = Inches(1.85 + row_i * 0.7)
    bg = DARK_CARD if row_i % 2 == 0 else RGBColor(0x15, 0x1F, 0x30)
    add_shape(slide, Inches(0.4), y, Inches(12.4), Inches(0.6), fill_color=bg)
    add_text_box(slide, col_x[0], Emu(y + Pt(4)), col_w[0], Inches(0.5), misconf, font_size=12, color=color, bold=True)
    add_text_box(slide, col_x[1], Emu(y + Pt(4)), col_w[1], Inches(0.5), impact, font_size=12, color=GRAY)
    add_text_box(slide, col_x[2], Emu(y + Pt(4)), col_w[2], Inches(0.5), fix, font_size=12, color=GREEN)


# ============================================================
# SLIDE 5: MITRE ATT&CK Mapping
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "MITRE ATT\u0026CK Kill Chain", font_size=32, color=PURPLE, bold=True)

techniques = [
    ("T1190", "Initial Access", "Exploit Public-Facing App", "Spring4Shell CVE-2022-22965", BLUE),
    ("T1059.004", "Execution", "Unix Shell", "Webshell command execution", BLUE),
    ("T1505.003", "Persistence", "Web Shell", "JSP webshell deployed", PURPLE),
    ("T1611", "Privilege Escalation", "Escape to Host", "nsenter + privileged pod", RED),
    ("T1552.007", "Credential Access", "Container API", "SA token theft, IMDS", ORANGE),
    ("T1613", "Discovery", "Container Discovery", "K8s resource enumeration", ORANGE),
    ("T1550.001", "Lateral Movement", "Application Access Token", "Cluster-admin token reuse", RED),
    ("T1530", "Collection", "Data from Cloud Storage", "AWS credentials exfiltration", RED),
]

for i, (tid, phase, name, evidence, color) in enumerate(techniques):
    y = Inches(1.2 + i * 0.72)
    # Phase box
    add_shape(slide, Inches(0.5), y, Inches(2.2), Inches(0.6), border_color=color)
    add_text_box(slide, Inches(0.6), Emu(y + Pt(2)), Inches(2), Inches(0.3), phase, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(0.6), Emu(y + Pt(20)), Inches(2), Inches(0.3), tid, font_size=10, color=GRAY)
    # Technique
    add_text_box(slide, Inches(2.9), Emu(y + Pt(6)), Inches(3.5), Inches(0.5), name, font_size=13, color=WHITE, bold=True)
    # Evidence
    add_text_box(slide, Inches(6.5), Emu(y + Pt(6)), Inches(6), Inches(0.5), evidence, font_size=12, color=GRAY)
    # Arrow
    if i < len(techniques) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.5), Emu(y + Inches(0.6)), Inches(0.2), Inches(0.12))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()


# ============================================================
# SLIDE 6: Cortex Response — Playbooks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Cortex Automated Response", font_size=32, color=GREEN, bold=True)

# Containment Playbook
add_shape(slide, Inches(0.3), Inches(1.2), Inches(4.0), Inches(5.8), border_color=GREEN)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(3.6), Inches(0.4),
             "Containment Playbook", font_size=16, color=GREEN, bold=True)
add_text_box(slide, Inches(0.5), Inches(1.75), Inches(3.6), Inches(0.3),
             "10 tasks — auto/manual gate", font_size=11, color=GRAY)

containment_steps = [
    "1. Triage (IOC extraction)",
    "2. Collect Evidence (Lambda)",
    "3. Severity Gate",
    "   Critical+Spring4Shell = auto",
    "4. Network Isolation (deny-all)",
    "5. Revoke RBAC (cluster-admin)",
    "6. Scale Down (replicas=0)",
    "7. Cordon Node",
    "8. Kill Pods (force delete)",
    "9. Verify Containment",
]
add_bullet_list(slide, Inches(0.5), Inches(2.2), Inches(3.6), Inches(4.5),
                containment_steps, font_size=11, bullet_color=WHITE)

# Forensic Playbook
add_shape(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(5.8), border_color=PURPLE)
add_text_box(slide, Inches(4.8), Inches(1.3), Inches(3.6), Inches(0.4),
             "Forensic Analysis Playbook", font_size=16, color=PURPLE, bold=True)
add_text_box(slide, Inches(4.8), Inches(1.75), Inches(3.6), Inches(0.3),
             "9 tasks — 5 XQL auto-exec", font_size=11, color=GRAY)

forensic_steps = [
    "1. Triage (IOC extraction)",
    "2. CVE + MITRE + XQL generation",
    "3. XQL: Causality Chain",
    "4. XQL: File Operations",
    "5. XQL: Network Connections",
    "6. XQL: Container Escape",
    "7. XQL: Credential Access",
    "8. Live Evidence (Lambda)",
]
add_bullet_list(slide, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.5),
                forensic_steps, font_size=11, bullet_color=WHITE)

# Search Similar Events Playbook
add_shape(slide, Inches(8.9), Inches(1.2), Inches(4.0), Inches(5.8), border_color=CYAN)
add_text_box(slide, Inches(9.1), Inches(1.3), Inches(3.6), Inches(0.4),
             "Search Similar Events", font_size=16, color=CYAN, bold=True)
add_text_box(slide, Inches(9.1), Inches(1.75), Inches(3.6), Inches(0.3),
             "8 tasks — 3 XQL auto-exec", font_size=11, color=GRAY)

search_steps = [
    "1. Triage (IOC extraction)",
    "2. Generate Hunt Queries",
    "3. XQL: Webshell Hunt (all nodes)",
    "4. XQL: Escape Hunt (all nodes)",
    "5. XQL: IMDS Theft Hunt",
    "6. Analyst Review",
    "7/8. Escalate or Close",
]
add_bullet_list(slide, Inches(9.1), Inches(2.2), Inches(3.6), Inches(4.5),
                search_steps, font_size=11, bullet_color=WHITE)


# ============================================================
# SLIDE 7: Automation Scripts
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Cortex Automation Scripts (4)", font_size=32, color=PURPLE, bold=True)

scripts = [
    ("ExtractK8sContainerEscapeIOCs", "Triage & IOC Extraction",
     "Analyzes XDR issue fields to extract IOCs (container ID, namespace, node FQDN, process SHA256). "
     "Determines severity (Critical/High/Medium/Low). Detects Spring4Shell + webshell patterns.",
     "K8sEscape.* context keys", BLUE),
    ("InvokeK8sContainmentLambda", "Lambda Containment (SigV4)",
     "Invokes AWS Lambda from Cortex XSIAM (GCP). Pure SigV4 signing (no boto3). "
     "STS AssumeRole + Lambda Invoke. 7 containment actions.",
     "K8sContainment.* context keys", ORANGE),
    ("K8sForensicAnalysis", "CVE / MITRE / XQL",
     "CVE enrichment (Spring4Shell CVSS 9.8). MITRE ATT&CK mapping (9 techniques). "
     "Generates 5 XQL forensic queries for auto-execution.",
     "K8sForensic.* context keys", PURPLE),
    ("K8sSearchSimilarEvents", "Cross-Tenant Threat Hunt",
     "Generates targeted + broad XQL queries. Lateral movement detection, "
     "blast radius assessment, webshell/escape/IMDS hunts.",
     "K8sSimilar.* context keys", CYAN),
]

for i, (name, subtitle, desc, output, color) in enumerate(scripts):
    y = Inches(1.1 + i * 1.5)
    add_shape(slide, Inches(0.4), y, Inches(12.4), Inches(1.35), border_color=color)
    add_text_box(slide, Inches(0.7), Emu(y + Pt(4)), Inches(5), Inches(0.4), name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(6), Emu(y + Pt(4)), Inches(4), Inches(0.4), subtitle, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.7), Emu(y + Pt(24)), Inches(11), Inches(0.5), desc, font_size=11, color=GRAY)
    add_text_box(slide, Inches(0.7), Emu(y + Pt(48)), Inches(11), Inches(0.3), f"Output: {output}", font_size=10, color=RGBColor(0x64, 0x74, 0x8B))


# ============================================================
# SLIDE 8: Lambda Containment Actions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Lambda Containment Actions", font_size=32, color=ORANGE, bold=True)

add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
             "AWS Lambda authenticates to EKS via STS presigned URL (x-k8s-aws-id)", font_size=14, color=GRAY)

actions = [
    ("collect_evidence", "Pod details, logs, events, RBAC audit, node status", GREEN),
    ("network_isolate", "Apply deny-all NetworkPolicy in namespace", ORANGE),
    ("revoke_rbac", "Delete cluster-admin ClusterRoleBinding", RED),
    ("scale_down", "Scale deployment to 0 replicas", ORANGE),
    ("cordon_node", "Mark node as unschedulable", ORANGE),
    ("delete_pod", "Force delete all pods in namespace", RED),
    ("full_containment", "Execute all actions in sequence", RED),
]

for i, (action, desc, color) in enumerate(actions):
    y = Inches(1.6 + i * 0.75)
    add_shape(slide, Inches(0.5), y, Inches(3.5), Inches(0.6), border_color=color)
    add_text_box(slide, Inches(0.7), Emu(y + Pt(8)), Inches(3.2), Inches(0.4), action, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(4.3), Emu(y + Pt(8)), Inches(8), Inches(0.4), desc, font_size=13, color=WHITE)


# ============================================================
# SLIDE 9: IAM Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "IAM Architecture", font_size=32, color=BLUE, bold=True)

# Dashboard flow
add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.5), border_color=BLUE)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
             "Dashboard (permanent credentials)", font_size=16, color=BLUE, bold=True)

txBox = add_text_box(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(1.5), "", font_size=13)
tf = txBox.text_frame
tf.paragraphs[0].text = "dashboard-user (Access Key)"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = WHITE
add_paragraph(tf, "    \u2193 AssumeRole", font_size=12, color=GRAY)
add_paragraph(tf, "dashboard-operator (scoped role)", font_size=13, color=WHITE)
add_paragraph(tf, "    EKS, ECR, Lambda, IAM, VPC, Logs", font_size=11, color=GRAY)

# Cortex flow
add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.5), border_color=PURPLE)
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5), Inches(0.4),
             "Cortex Playbook (permanent credentials)", font_size=16, color=PURPLE, bold=True)

txBox = add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(1.5), "", font_size=13)
tf = txBox.text_frame
tf.paragraphs[0].text = "cortex-playbook-user (Access Key)"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = WHITE
add_paragraph(tf, "    \u2193 AssumeRole", font_size=12, color=GRAY)
add_paragraph(tf, "lambda-invoker (scoped role)", font_size=13, color=WHITE)
add_paragraph(tf, "    lambda:InvokeFunction only", font_size=11, color=GRAY)

# Bootstrap note
add_shape(slide, Inches(0.3), Inches(4.2), Inches(12.7), Inches(1.5), border_color=GRAY)
add_text_box(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
             "Bootstrap (one-time, admin credentials)", font_size=16, color=GRAY, bold=True)
add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.8),
             "Admin credentials \u2192 terraform apply (terraform-infra/ + terraform-lambda/)\n"
             "Creates all IAM users, roles, EKS, Lambda. Admin credentials no longer needed after setup.",
             font_size=12, color=GRAY)


# ============================================================
# SLIDE 10: Demo Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Demo Flow", font_size=32, color=GREEN, bold=True)

demo_steps = [
    ("1", "Configure AWS", "Enter admin credentials in dashboard", BLUE),
    ("2", "Deploy Infrastructure", "terraform apply \u2192 VPC, EKS, ECR, IAM (~15 min)", BLUE),
    ("3", "Switch to Dashboard User", "Permanent credentials from terraform output", BLUE),
    ("4", "Build & Push Image", "Docker buildx (linux/amd64) + push to ECR", CYAN),
    ("5", "Deploy App", "K8s manifests: privileged pod + LoadBalancer", CYAN),
    ("6", "Run Attack Chain", "Step 1: RCE \u2192 Step 2: Escape \u2192 Step 3: Takeover", RED),
    ("7", "Deploy Cortex", "Push 4 scripts + 3 playbooks + Lambda", PURPLE),
    ("8", "Observe Detection", "XDR creates issue \u2192 triggers playbook", GREEN),
    ("9", "Automated Response", "Containment: NetworkPolicy, RBAC, scale, cordon, kill", GREEN),
    ("10", "Cleanup", "Destroy Lambda + Destroy All", GRAY),
]

for i, (num, title, desc, color) in enumerate(demo_steps):
    y = Inches(1.1 + i * 0.6)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, Inches(0.5), Emu(y + Pt(4)), Inches(0.45), Inches(0.4), num, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title + desc
    add_text_box(slide, Inches(1.2), Emu(y + Pt(2)), Inches(3.5), Inches(0.4), title, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(4.8), Emu(y + Pt(4)), Inches(8), Inches(0.4), desc, font_size=13, color=GRAY)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "K8s_Container_Escape_Demo.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
