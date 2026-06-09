"""
Generate the Code-to-Cloud-to-SOC pitch deck using the Cortex 2025 template.
The deck inherits theme, slide masters, fonts and colors from CFTD_10062026.pptx.

Run:
  python3 docs/build/build_pitch.py
Output:
  docs/pitch.pptx
"""

import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ------------------------------------------------------------------ paths
ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = Path("/Users/cley/Downloads/CFTD_10062026.pptx")
OUTPUT = ROOT / "docs" / "pitch.pptx"

# ------------------------------------------------------------------ brand
FONT = "Helvetica Neue"
CORTEX_ORANGE = RGBColor(0xFA, 0x58, 0x2D)
CORTEX_ORANGE_DARK = RGBColor(0xB2, 0x38, 0x08)
CORTEX_YELLOW = RGBColor(0xFF, 0xCB, 0x06)
CORTEX_CYAN = RGBColor(0x00, 0xC0, 0xE8)
CORTEX_GREEN = RGBColor(0x00, 0xCC, 0x66)
CORTEX_TEAL = RGBColor(0x00, 0xAE, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHT_GREY = RGBColor(0xE9, 0xE9, 0xE9)

# Slide canvas is 20.0in x 11.25in
SLIDE_W = Inches(20)
SLIDE_H = Inches(11.25)


# ------------------------------------------------------------------ helpers
def delete_all_slides(prs):
    """Remove every existing slide while preserving masters/layouts/theme."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for sldId in slides:
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def find_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(f"layout {name!r} not found")


def set_run(run, text, *, size=18, bold=False, color=WHITE, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def clear_tf(tf):
    """Remove every paragraph and reset to a single empty one."""
    p_elements = tf._txBody.findall(qn("a:p"))
    for p in p_elements[1:]:
        tf._txBody.remove(p)
    first = p_elements[0]
    for r in first.findall(qn("a:r")):
        first.remove(r)
    for br in first.findall(qn("a:br")):
        first.remove(br)
    for fld in first.findall(qn("a:fld")):
        first.remove(fld)


def add_paragraph(tf, lines, *, size=18, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, bullet=False, space_after=6):
    """Append a paragraph (or several) to a text frame."""
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        if tf.paragraphs and i == 0 and not tf.paragraphs[-1].runs and not tf.paragraphs[-1].text:
            p = tf.paragraphs[-1]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if bullet:
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "•")
            buFont = etree.SubElement(pPr, qn("a:buFont"))
            buFont.set("typeface", FONT)
            pPr.set("indent", "-228600")
            pPr.set("marL", "228600")
        run = p.add_run()
        set_run(run, line, size=size, bold=bold, color=color)


def add_text_box(slide, left, top, width, height, *, fill=None, line=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.word_wrap = True
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(1.5)
    else:
        box.line.fill.background()
    clear_tf(tf)
    return box, tf


def add_accent_bar(slide, left, top, width, height, color):
    """Add a solid colored rectangle as an accent bar."""
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar.text_frame.text = ""
    return bar


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    clear_tf(notes_tf)
    for i, paragraph in enumerate(text.strip().split("\n\n")):
        if i == 0:
            p = notes_tf.paragraphs[0]
        else:
            p = notes_tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = paragraph.strip()
        run.font.name = FONT
        run.font.size = Pt(12)


def set_slide_bg(slide, color):
    """Force the slide background to a solid color (overrides master)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ------------------------------------------------------------------ slide builders
def build_cover(prs):
    layout = find_layout(prs, "TITLE_1_2_3_1_1")
    s = prs.slides.add_slide(layout)
    set_slide_bg(s, BLACK)

    # Decorative orange bar
    add_accent_bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.35), CORTEX_ORANGE)

    # Title
    _, tf = add_text_box(s, Inches(1.0), Inches(3.5), Inches(18), Inches(2.5))
    add_paragraph(tf, "From Code to Cloud to SOC",
                  size=72, bold=True, color=WHITE)

    # Subtitle
    _, tf2 = add_text_box(s, Inches(1.0), Inches(5.8), Inches(18), Inches(1.2))
    add_paragraph(tf2, "See the attack. Stop the breach. Automatically.",
                  size=36, bold=False, color=CORTEX_ORANGE)

    # Tagline
    _, tf3 = add_text_box(s, Inches(1.0), Inches(7.4), Inches(18), Inches(1.5))
    add_paragraph(tf3,
                  "A live demonstration of how Cortex Cloud detects and contains",
                  size=22, color=LIGHT_GREY)
    add_paragraph(tf3,
                  "a cloud-native attack — end to end, in under a minute.",
                  size=22, color=LIGHT_GREY)

    # Footer band
    add_accent_bar(s, Inches(0), Inches(10.85), SLIDE_W, Inches(0.4), CORTEX_ORANGE_DARK)
    _, tff = add_text_box(s, Inches(1.0), Inches(10.4), Inches(18), Inches(0.45))
    add_paragraph(tff,
                  "Palo Alto Networks  ·  Partner enablement  ·  10–15 min",
                  size=14, color=LIGHT_GREY)

    add_notes(s, """SCRIPT FR (1 min) :
"Bonjour à tous. Aujourd'hui je vais vous montrer comment Cortex Cloud permet à vos clients de voir une attaque cloud-native se dérouler en temps réel — et surtout, c'est ça le point clé, de la stopper automatiquement, sans intervention humaine, en moins d'une minute.

Cette démo, vous allez pouvoir la rejouer chez vos clients en 15 minutes chrono. C'est un outil commercial puissant."
""")
    return s


def slide_with_header(prs, title, eyebrow=None):
    """Blank slide with the standard Cortex 2025 title bar."""
    layout = find_layout(prs, "BLANK_4") if any(
        l.name == "BLANK_4" for m in prs.slide_masters for l in m.slide_layouts
    ) else find_layout(prs, "TITLE_AND_BODY_1_1")
    s = prs.slides.add_slide(layout)
    set_slide_bg(s, BLACK)

    # Remove inherited placeholders for full control
    for ph in list(s.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Top accent bar
    add_accent_bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.25), CORTEX_ORANGE)

    # Eyebrow
    if eyebrow:
        _, tfe = add_text_box(s, Inches(0.8), Inches(0.5), Inches(18), Inches(0.5))
        add_paragraph(tfe, eyebrow.upper(),
                      size=16, bold=True, color=CORTEX_ORANGE)

    # Title
    top = Inches(1.0) if eyebrow else Inches(0.7)
    _, tft = add_text_box(s, Inches(0.8), top, Inches(18), Inches(1.4))
    add_paragraph(tft, title, size=44, bold=True, color=WHITE)

    # Footer
    add_accent_bar(s, Inches(0), Inches(10.95), SLIDE_W, Inches(0.3), CORTEX_ORANGE_DARK)
    _, tff = add_text_box(s, Inches(0.8), Inches(10.55), Inches(18), Inches(0.4))
    add_paragraph(tff,
                  "Cortex Cloud  ·  Code to Cloud to SOC  ·  Palo Alto Networks",
                  size=12, color=LIGHT_GREY)
    return s


# ------------------------------------------------------------------ KPI helper
def add_kpi_card(slide, left, top, width, height, value, label, color=CORTEX_ORANGE):
    # Card background
    from pptx.enum.shapes import MSO_SHAPE
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x1F, 0x2E)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    card.shadow.inherit = False
    card.text_frame.text = ""

    # Value
    _, tf = add_text_box(slide, left, top + Inches(0.6), width, Inches(2.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, value, size=88, bold=True, color=color)

    # Label
    _, tf2 = add_text_box(slide, left, top + height - Inches(2.0), width, Inches(1.8))
    for line in label.split("\n"):
        p = tf2.add_paragraph() if tf2.paragraphs[0].runs else tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        set_run(run, line, size=18, color=LIGHT_GREY)


# ------------------------------------------------------------------ deck
def build_deck():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")

    # Copy template to output to keep theme/masters/embedded fonts
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(TEMPLATE, OUTPUT)

    prs = Presentation(OUTPUT)
    delete_all_slides(prs)

    # ====================== Slide 1 — Cover
    build_cover(prs)

    # ====================== Slide 2 — Cloud risk explodes
    s = slide_with_header(prs, "The cloud risk is exploding", eyebrow="Why we are here")

    kpi_top = Inches(3.2)
    kpi_w = Inches(5.5)
    kpi_h = Inches(4.2)
    gap = Inches(0.7)
    total_w = 3 * kpi_w + 2 * gap
    start_left = (SLIDE_W - total_w) / 2

    add_kpi_card(s, start_left, kpi_top, kpi_w, kpi_h,
                 "+95%", "of organizations experienced\na cloud security incident",
                 color=CORTEX_ORANGE)
    add_kpi_card(s, start_left + kpi_w + gap, kpi_top, kpi_w, kpi_h,
                 "277", "days average to detect\nand contain a breach",
                 color=CORTEX_YELLOW)
    add_kpi_card(s, start_left + 2 * (kpi_w + gap), kpi_top, kpi_w, kpi_h,
                 "$4.9M", "average cost of a\ncloud data breach",
                 color=CORTEX_CYAN)

    _, tfq = add_text_box(s, Inches(1.0), Inches(8.2), Inches(18), Inches(1.8),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "The attack surface has moved.",
                  size=24, bold=True, color=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "Workloads, identities, configurations, secrets — all in the cloud, all interconnected, all exploitable.",
                  size=20, color=LIGHT_GREY)

    add_notes(s, """SCRIPT FR (1 min30) :
"Les chiffres parlent d'eux-mêmes : 95% des organisations ont subi un incident cloud sur les 12 derniers mois. Le temps moyen pour détecter ET contenir une brèche reste de 277 jours — neuf mois. Et le coût moyen d'une brèche cloud atteint 4,9 millions de dollars.

La réalité : l'attaque ne vient plus du périmètre. Elle vient de l'intérieur — d'une mauvaise config, d'une identité compromise, d'un container qui s'échappe. Vos clients le savent. Ils cherchent une réponse."
""")

    # ====================== Slide 3 — Why customers can't test
    s = slide_with_header(prs, "Why customers can't test their own defenses",
                          eyebrow="The customer pain")

    col_top = Inches(2.8)
    col_w = Inches(8.8)
    col_h = Inches(6.2)
    col_gap = Inches(0.4)
    left_x = Inches(0.8)
    right_x = left_x + col_w + col_gap

    _, tfL = add_text_box(s, left_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_ORANGE)
    add_paragraph(tfL, "It's hard", size=28, bold=True, color=CORTEX_ORANGE,
                  space_after=12)
    add_paragraph(tfL,
                  "Building a realistic attack scenario requires deep offensive skills.",
                  size=20, color=WHITE, bullet=True, space_after=10)
    add_paragraph(tfL,
                  "Most security teams are defenders, not red teamers.",
                  size=20, color=WHITE, bullet=True, space_after=10)
    add_paragraph(tfL,
                  "Engaging a pentest firm = weeks of planning, significant cost.",
                  size=20, color=WHITE, bullet=True)

    _, tfR = add_text_box(s, right_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_ORANGE)
    add_paragraph(tfR, "It's risky", size=28, bold=True, color=CORTEX_ORANGE,
                  space_after=12)
    add_paragraph(tfR,
                  "Testing in production = business disruption.",
                  size=20, color=WHITE, bullet=True, space_after=10)
    add_paragraph(tfR,
                  "Building a lab from scratch = months of effort.",
                  size=20, color=WHITE, bullet=True, space_after=10)
    add_paragraph(tfR,
                  "Tools are fragmented across CSPM, CWPP, XDR, SOAR…",
                  size=20, color=WHITE, bullet=True)

    _, tfq = add_text_box(s, Inches(0.8), Inches(9.3), Inches(18.4), Inches(1.4),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "Result: customers buy security tools but never validate they actually detect and stop an attack — until it's too late.",
                  size=20, bold=True, color=CORTEX_ORANGE)

    add_notes(s, """SCRIPT FR (2 min) :
"Voilà le vrai problème de vos clients : ils achètent des outils de sécurité, mais ils ne testent jamais s'ils marchent vraiment.

Pourquoi ? Parce que c'est dur — il faut des compétences offensives que les équipes défensives n'ont pas. Parce que c'est risqué — on ne joue pas une attaque en prod. Parce que monter un lab prend des mois. Et parce qu'avec 5 ou 6 outils empilés, personne ne sait qui détecte quoi.

Conclusion : ils découvrent que leur stack ne marche pas... le jour de la vraie attaque. C'est exactement ce que Cortex Cloud résout."
""")

    # ====================== Slide 4 — One platform, three moments
    s = slide_with_header(prs, "One platform, three moments",
                          eyebrow="The Cortex Cloud answer")

    card_top = Inches(2.8)
    card_w = Inches(5.8)
    card_h = Inches(5.8)
    card_gap = Inches(0.5)
    cards_total = 3 * card_w + 2 * card_gap
    start = (SLIDE_W - cards_total) / 2

    cards = [
        ("BEFORE", "Shift-Left", CORTEX_CYAN,
         "Find vulnerabilities and misconfigurations in the code,",
         "before they reach production."),
        ("DURING", "Runtime Protection", CORTEX_ORANGE,
         "See and stop attacks as they happen,",
         "on every workload, every identity, every cloud."),
        ("AFTER", "Automated Response", CORTEX_GREEN,
         "Contain, investigate and remediate without",
         "human intervention — in seconds."),
    ]

    for i, (eyebrow, title, color, line1, line2) in enumerate(cards):
        x = start + i * (card_w + card_gap)
        from pptx.enum.shapes import MSO_SHAPE
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x24)
        card.line.color.rgb = color
        card.line.width = Pt(2.5)
        card.shadow.inherit = False
        card.text_frame.text = ""

        _, tfe = add_text_box(s, x, card_top + Inches(0.4), card_w, Inches(0.5))
        p = tfe.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); set_run(run, eyebrow, size=20, bold=True, color=color)

        _, tft = add_text_box(s, x, card_top + Inches(1.2), card_w, Inches(1.0))
        p = tft.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); set_run(run, title, size=32, bold=True, color=WHITE)

        _, tfb = add_text_box(s, x + Inches(0.3), card_top + Inches(3.0),
                              card_w - Inches(0.6), Inches(2.5))
        for line in (line1, line2):
            p = tfb.add_paragraph() if tfb.paragraphs[0].runs else tfb.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); set_run(run, line, size=18, color=LIGHT_GREY)

    _, tfq = add_text_box(s, Inches(0.8), Inches(9.3), Inches(18.4), Inches(1.4),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "The unique value: one platform that sees the code, the cloud, and the SOC. No silo. No blind spot.",
                  size=20, bold=True, color=CORTEX_ORANGE)

    add_notes(s, """SCRIPT FR (1 min30) :
"Cortex Cloud, c'est trois moments, une seule plateforme.
Avant l'attaque : on trouve les vulnérabilités dans le code, avant qu'elles n'arrivent en prod. C'est le shift-left.
Pendant l'attaque : on voit ce qui se passe en temps réel, sur chaque workload, chaque identité, chaque cloud.
Après — ou plutôt pendant, parce qu'on parle de quelques secondes : on contient et on remédie automatiquement.

La différence avec la concurrence : un seul agent, une seule console, une seule donnée partagée entre les trois moments. Pas de silo."
""")

    # ====================== Slide 5 — The demo
    s = slide_with_header(prs, "The demo — what your customer will see",
                          eyebrow="The wow moment")

    col_top = Inches(2.8)
    col_w = Inches(8.8)
    col_h = Inches(6.2)
    col_gap = Inches(0.4)
    left_x = Inches(0.8); right_x = left_x + col_w + col_gap

    _, tfL = add_text_box(s, left_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_ORANGE)
    add_paragraph(tfL, "A realistic attack chain", size=26, bold=True,
                  color=CORTEX_ORANGE, space_after=12)
    for b in [
        "Starts from a vulnerable web app in production",
        "Progresses to container escape",
        "Ends with full cluster takeover and data theft",
        "≈ 2 minutes, fully reproducible",
    ]:
        add_paragraph(tfL, b, size=20, color=WHITE, bullet=True, space_after=10)

    _, tfR = add_text_box(s, right_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_GREEN)
    add_paragraph(tfR, "The platform reaction", size=26, bold=True,
                  color=CORTEX_GREEN, space_after=12)
    for b in [
        "Real-time alerts with full context",
        "Automatic correlation across code + cloud + endpoint",
        "Playbooks fire without analyst input",
        "Workload contained in seconds",
    ]:
        add_paragraph(tfR, b, size=20, color=WHITE, bullet=True, space_after=10)

    _, tfq = add_text_box(s, Inches(0.8), Inches(9.3), Inches(18.4), Inches(1.4),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "One click in the dashboard → the attack runs → Cortex Cloud reacts → the customer sees the full story.",
                  size=20, bold=True, color=CORTEX_ORANGE)

    add_notes(s, """SCRIPT FR (1 min30) :
"La démo en elle-même : on lance une attaque réaliste — pas un truc artificiel, un scénario que vos clients vivent vraiment — qui part d'une appli web vulnérable, qui s'échappe du conteneur, qui prend le contrôle du cluster.

En face, Cortex Cloud réagit en temps réel : alertes contextualisées, corrélation automatique entre les couches, playbooks qui se déclenchent seuls, et workload coupé en quelques secondes.

Tout ça en un clic dans le dashboard. C'est le 'wow moment'."
""")

    # ====================== Slide 6 — Detection
    s = slide_with_header(prs, "Detection — Cortex Cloud sees everything",
                          eyebrow="Full visibility")

    col_top = Inches(2.8)
    col_w = Inches(8.8); col_h = Inches(6.2); col_gap = Inches(0.4)
    left_x = Inches(0.8); right_x = left_x + col_w + col_gap

    _, tfL = add_text_box(s, left_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_CYAN)
    add_paragraph(tfL, "Full visibility across the kill chain", size=24, bold=True,
                  color=CORTEX_CYAN, space_after=12)
    for b in [
        "Code vulnerabilities flagged before deployment",
        "Misconfigurations detected in cloud posture",
        "Runtime behavior monitored on every workload",
        "Identity abuse tracked across cloud accounts",
        "Lateral movement mapped in real time",
    ]:
        add_paragraph(tfL, b, size=19, color=WHITE, bullet=True, space_after=8)

    _, tfR = add_text_box(s, right_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_CYAN)
    add_paragraph(tfR, "Unified telemetry", size=24, bold=True,
                  color=CORTEX_CYAN, space_after=12)
    for b in [
        "Single console, single timeline",
        "Native MITRE ATT&CK mapping",
        "Auto-correlation across cloud, container, identity, endpoint",
        "Full incident story, already assembled",
    ]:
        add_paragraph(tfR, b, size=19, color=WHITE, bullet=True, space_after=8)

    _, tfq = add_text_box(s, Inches(0.8), Inches(9.3), Inches(18.4), Inches(1.4),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "The analyst doesn't hunt for evidence — Cortex delivers the full incident story, already assembled.",
                  size=20, bold=True, color=CORTEX_ORANGE)

    add_notes(s, """SCRIPT FR (2 min) :
"Côté détection, Cortex Cloud voit toute la chaîne d'attaque, pas juste un bout. Les vulnérabilités côté code, les mauvaises configs côté cloud, le comportement runtime sur chaque workload, l'abus d'identité entre comptes AWS, et le mouvement latéral en temps réel.

Tout converge dans une seule console, avec une seule timeline, mappée nativement sur MITRE ATT&CK. Le résultat : l'analyste ne perd plus des heures à reconstituer ce qui s'est passé. Cortex lui livre l'histoire complète, déjà assemblée.

C'est ça qui change la vie d'un SOC."
""")

    # ====================== Slide 7 — Response (automated)
    s = slide_with_header(prs, "Response — Cortex Cloud acts, automatically",
                          eyebrow="From alert to containment, hands-off")

    # Top: workflow steps as a horizontal pipeline
    steps = [
        ("1", "Detection", "Multi-source correlation"),
        ("2", "Enrichment", "Forensic context auto-gathered"),
        ("3", "Decision", "Playbook applies policy"),
        ("4", "Action", "Workload contained, identity revoked"),
    ]
    step_top = Inches(2.6)
    step_w = Inches(4.3); step_h = Inches(3.2); step_gap = Inches(0.35)
    total = 4 * step_w + 3 * step_gap
    start = (SLIDE_W - total) / 2

    from pptx.enum.shapes import MSO_SHAPE
    for i, (num, title, desc) in enumerate(steps):
        x = start + i * (step_w + step_gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, step_top, step_w, step_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x24)
        card.line.color.rgb = CORTEX_GREEN
        card.line.width = Pt(2)
        card.shadow.inherit = False
        card.text_frame.text = ""

        _, tfn = add_text_box(s, x, step_top + Inches(0.3), step_w, Inches(0.9))
        p = tfn.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); set_run(run, num, size=48, bold=True, color=CORTEX_GREEN)

        _, tft = add_text_box(s, x, step_top + Inches(1.3), step_w, Inches(0.7))
        p = tft.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); set_run(run, title, size=22, bold=True, color=WHITE)

        _, tfd = add_text_box(s, x + Inches(0.2), step_top + Inches(2.1),
                              step_w - Inches(0.4), Inches(1.0))
        p = tfd.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); set_run(run, desc, size=15, color=LIGHT_GREY)

    # Bottom KPIs
    kpi_top = Inches(6.4)
    kpi_w = Inches(5.5); kpi_h = Inches(2.6); kpi_gap = Inches(0.6)
    total_w = 3 * kpi_w + 2 * kpi_gap
    start = (SLIDE_W - total_w) / 2

    add_kpi_card(s, start, kpi_top, kpi_w, kpi_h, "< 60s", "MTTR (automated)",
                 color=CORTEX_ORANGE)
    add_kpi_card(s, start + kpi_w + kpi_gap, kpi_top, kpi_w, kpi_h,
                 "0", "Analyst clicks required", color=CORTEX_YELLOW)
    add_kpi_card(s, start + 2 * (kpi_w + kpi_gap), kpi_top, kpi_w, kpi_h,
                 "24/7", "Consistent response", color=CORTEX_GREEN)

    _, tfq = add_text_box(s, Inches(0.8), Inches(9.3), Inches(18.4), Inches(1.4),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "Detection alone is table stakes. Automated remediation is the differentiator.",
                  size=20, bold=True, color=CORTEX_ORANGE)

    add_notes(s, """SCRIPT FR (2 min) :
"Et c'est là le vrai différentiateur de Cortex Cloud : on ne se contente pas de détecter, on RÉPOND. Automatiquement.

Le workflow : détection multi-source, enrichissement forensique automatique, décision via playbook, action immédiate. Le workload compromis est isolé, le process malicieux tué, les tokens volés révoqués, l'image vulnérable bloquée au redéploiement.

Tout ça en moins de 60 secondes, sans intervention humaine, 24/7, de manière 100% reproductible. Comparez ça aux 277 jours de moyenne du marché. Le ROI est immédiat."
""")

    # ====================== Slide 8 — Before / After
    s = slide_with_header(prs, "Before / After — the customer outcome",
                          eyebrow="Measurable impact")

    # Build a manual two-column comparison table
    rows = [
        ("",                          "Before Cortex Cloud",            "After Cortex Cloud"),
        ("Vulnerabilities",           "Discovered in production",       "Caught before deployment"),
        ("Attack visibility",         "Fragmented across 5+ tools",     "Single pane, full context"),
        ("Detection time",            "Hours to days",                  "Seconds"),
        ("Response",                  "Manual, inconsistent",           "Automated, 24/7"),
        ("MTTR",                      "Days",                           "Under 1 minute"),
        ("Compliance",                "Annual audit panic",             "Continuous posture"),
    ]

    tbl_left = Inches(0.8)
    tbl_top = Inches(2.6)
    tbl_w = Inches(18.4)
    tbl_h = Inches(6.2)
    table = s.shapes.add_table(len(rows), 3, tbl_left, tbl_top, tbl_w, tbl_h).table
    table.columns[0].width = Inches(4.4)
    table.columns[1].width = Inches(7.0)
    table.columns[2].width = Inches(7.0)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = CORTEX_ORANGE
                font_color = BLACK
                bold = True
                size = 18
            else:
                cell.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x24) if ri % 2 == 1 else RGBColor(0x0A, 0x0E, 0x1A)
                if ci == 0:
                    font_color = CORTEX_ORANGE; bold = True; size = 16
                elif ci == 1:
                    font_color = LIGHT_GREY; bold = False; size = 16
                else:
                    font_color = WHITE; bold = True; size = 16
            tf = cell.text_frame
            clear_tf(tf)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ri == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            set_run(run, val, size=size, bold=bold, color=font_color)
            cell.margin_left = Inches(0.2); cell.margin_right = Inches(0.2)
            cell.margin_top = Inches(0.1); cell.margin_bottom = Inches(0.1)

    _, tfq = add_text_box(s, Inches(0.8), Inches(9.3), Inches(18.4), Inches(1.4),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "The shift: from reactive and manual to proactive and autonomous.",
                  size=22, bold=True, color=CORTEX_ORANGE)

    add_notes(s, """SCRIPT FR (1 min30) :
"Ce tableau est celui que vous devez laisser à votre interlocuteur.

Avant Cortex : vulnérabilités en prod, visibilité fragmentée, détection en heures ou jours, réponse manuelle, MTTR en jours, compliance en mode panique annuelle.
Après Cortex : tout est inversé. Détection en secondes, réponse automatisée, MTTR en moins d'une minute, posture continue.

Le message à retenir : on passe d'une sécurité réactive et manuelle à une sécurité proactive et autonome. C'est la promesse, et c'est exactement ce que la démo prouve."
""")

    # ====================== Slide 9 — Why partners win
    s = slide_with_header(prs, "Why this matters for you, our partners",
                          eyebrow="Partner value")

    col_top = Inches(2.8)
    col_w = Inches(8.8); col_h = Inches(6.2); col_gap = Inches(0.4)
    left_x = Inches(0.8); right_x = left_x + col_w + col_gap

    _, tfL = add_text_box(s, left_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_YELLOW)
    add_paragraph(tfL, "A repeatable sales asset", size=26, bold=True,
                  color=CORTEX_YELLOW, space_after=12)
    for b in [
        "One-click demo at every customer meeting",
        "15 minutes to deliver the wow moment",
        "Works on any Kubernetes (AWS, on-prem, customer's own)",
        "Zero install required at the customer",
    ]:
        add_paragraph(tfL, b, size=19, color=WHITE, bullet=True, space_after=10)

    _, tfR = add_text_box(s, right_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_ORANGE)
    add_paragraph(tfR, "A platform story that wins", size=26, bold=True,
                  color=CORTEX_ORANGE, space_after=12)
    for b in [
        "Competitors show one layer — we show the full chain",
        "Justifies the CNAPP + XDR + SOAR consolidation",
        "Opens the door to larger deals and strategic accounts",
        "Differentiates you from product-only resellers",
    ]:
        add_paragraph(tfR, b, size=19, color=WHITE, bullet=True, space_after=10)

    _, tfq = add_text_box(s, Inches(0.8), Inches(9.3), Inches(18.4), Inches(1.4),
                          line=CORTEX_ORANGE)
    add_paragraph(tfq,
                  "You don't sell a tool. You sell an outcome — and you can prove it in 15 minutes.",
                  size=22, bold=True, color=CORTEX_ORANGE)

    add_notes(s, """SCRIPT FR (1 min30) :
"Pourquoi cette démo est un game-changer pour vous :

Côté commercial — c'est un asset rejouable en un clic chez n'importe quel client, en 15 minutes, sur n'importe quel Kubernetes. Pas d'installation, pas de prérequis lourds.

Côté stratégique — vos concurrents montrent une seule couche, vous montrez toute la chaîne. C'est ce qui justifie la consolidation CNAPP + XDR + SOAR, c'est ce qui ouvre les gros deals, c'est ce qui vous différencie.

Le message : vous ne vendez plus un outil, vous vendez un résultat — et vous pouvez le prouver en 15 minutes devant le client."
""")

    # ====================== Slide 10 — Next steps
    s = slide_with_header(prs, "Next steps", eyebrow="Let's get started")

    col_top = Inches(2.8)
    col_w = Inches(8.8); col_h = Inches(5.0); col_gap = Inches(0.4)
    left_x = Inches(0.8); right_x = left_x + col_w + col_gap

    _, tfL = add_text_box(s, left_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_GREEN)
    add_paragraph(tfL, "Book a customer POV", size=26, bold=True,
                  color=CORTEX_GREEN, space_after=12)
    for b in [
        "1–2 hours of customer time",
        "We provide the environment",
        "Output: live demo + tailored posture report",
    ]:
        add_paragraph(tfL, b, size=20, color=WHITE, bullet=True, space_after=10)

    _, tfR = add_text_box(s, right_x, col_top, col_w, col_h,
                          fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_CYAN)
    add_paragraph(tfR, "Enable your team", size=26, bold=True,
                  color=CORTEX_CYAN, space_after=12)
    for b in [
        "Train your pre-sales on the demo flow",
        "Co-selling support from Palo Alto channel SE",
        "Access to the partner enablement repo",
    ]:
        add_paragraph(tfR, b, size=20, color=WHITE, bullet=True, space_after=10)

    _, tfR2 = add_text_box(s, Inches(0.8), Inches(8.0), Inches(18.4), Inches(2.0),
                           fill=RGBColor(0x14, 0x18, 0x24), line=CORTEX_ORANGE)
    add_paragraph(tfR2, "Resources available today", size=22, bold=True,
                  color=CORTEX_ORANGE, space_after=8)
    add_paragraph(tfR2,
                  "Demo video  ·  Architecture diagrams  ·  Pitch deck  ·  Step-by-step partner guide",
                  size=18, color=LIGHT_GREY)
    add_paragraph(tfR2,
                  "Contact your Palo Alto channel SE to get started.",
                  size=18, color=WHITE, bold=True)

    add_notes(s, """SCRIPT FR (1 min) :
"Les prochaines étapes, très concrètement :

Pour un POV client : on a besoin d'une à deux heures de leur temps, on fournit l'environnement, et le client repart avec une démo live et un rapport personnalisé sur sa posture.

Pour votre équipe : on forme vos pré-sales sur le flow de démo, on vous accompagne en co-selling, et on vous donne accès au repo d'enablement partenaire.

Toutes les ressources sont prêtes — vidéo, architecture, ce deck, guide pas-à-pas. Contactez votre channel SE Palo Alto et on lance."
""")

    # ====================== Slide 11 — Questions
    s = slide_with_header(prs, "", eyebrow=None)
    # Centered big title
    _, tf = add_text_box(s, Inches(0.8), Inches(4.0), Inches(18.4), Inches(2.0))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); set_run(run, "Questions?", size=120, bold=True, color=CORTEX_ORANGE)

    _, tf2 = add_text_box(s, Inches(0.8), Inches(6.5), Inches(18.4), Inches(1.2))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, "Let's turn cloud risk into competitive advantage.",
            size=32, color=WHITE)

    _, tf3 = add_text_box(s, Inches(0.8), Inches(8.5), Inches(18.4), Inches(1.0))
    p = tf3.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, "Contact your Palo Alto channel SE",
            size=22, color=LIGHT_GREY)

    add_notes(s, """SCRIPT FR — Q&A, anticipez ces questions :

1. "Combien de temps pour déployer chez un client ?"
   → Environnement de démo prêt en 20-30 min. POV avec données client : 1 à 2 jours selon scope.

2. "Sur quels clouds ça marche ?"
   → AWS, GCP, Azure, et n'importe quel Kubernetes (on-prem inclus). La démo couvre AWS et on-prem aujourd'hui.

3. "Quelle différence vs Wiz, CrowdStrike, Sysdig ?"
   → Wiz : posture cloud uniquement, pas de runtime ni de SOC.
   → CrowdStrike : excellent endpoint, faible côté cloud-native.
   → Sysdig : runtime correct, pas de plateforme SOC unifiée.
   → Cortex Cloud : la seule plateforme qui couvre code + cloud + SOC avec un seul agent.

4. "Le client doit avoir tout l'écosystème Cortex pour en bénéficier ?"
   → Non. On peut commencer par une brique (CNAPP ou XDR) et étendre. Mais la vraie valeur — la chaîne complète — vient de la plateforme unifiée. C'est l'objectif de l'upsell.

5. "Quel est le pricing ?"
   → Modèle par workload + endpoint. Renvoyez vers votre channel manager pour un devis personnalisé.

6. "Combien de temps pour qu'un client soit autonome dessus ?"
   → Onboarding standard : 2 à 4 semaines avec accompagnement Pro Services ou partenaire certifié.
""")

    # ------------------------------------------------------------------ save
    prs.save(OUTPUT)
    print(f"✓ Generated {OUTPUT}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
